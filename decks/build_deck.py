#!/usr/bin/env python3
"""Revive seed deck v8. Real layout design: per-slide structure, scale contrast, dark moments,
ghost numerals + registration ticks as signature elements. Brand: paper/ink/cobalt, mono labels."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

PAPER      = RGBColor(0xf4,0xf5,0xf1)
PANEL      = RGBColor(0xfb,0xfc,0xf8)
PANEL2     = RGBColor(0xee,0xf0,0xeb)
GHOST      = RGBColor(0xe4,0xe7,0xdf)   # ghost numeral on paper
GHOST_D    = RGBColor(0x1e,0x24,0x31)   # ghost numeral on ink
INK        = RGBColor(0x15,0x19,0x22)
INK2       = RGBColor(0x1b,0x20,0x2b)
COBALT     = RGBColor(0x49,0x67,0xf2)
COBALT_DEEP= RGBColor(0x2e,0x49,0xc8)
COBALT_WASH= RGBColor(0xed,0xf0,0xff)
MUTED      = RGBColor(0x66,0x70,0x7e)
FAINT      = RGBColor(0x8a,0x92,0x9d)
HAIR       = RGBColor(0xc9,0xce,0xc8)
RED        = RGBColor(0xc2,0x41,0x3a)
RED_BR     = RGBColor(0xe8,0x6a,0x5f)
RED_WASH   = RGBColor(0xfc,0xed,0xeb)
GREEN      = RGBColor(0x14,0x80,0x60)
GREEN_BR   = RGBColor(0x3f,0xbf,0x8f)
GREEN_WASH = RGBColor(0xe9,0xf4,0xef)
SHADOW     = RGBColor(0xd9,0xdd,0xd6)
WHITE      = RGBColor(0xff,0xff,0xff)
LILAC      = RGBColor(0x8e,0xa0,0xff)
SLATE      = RGBColor(0x9a,0xa6,0xc8)

SANS = "Helvetica Neue"
MONO = "Menlo"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
I = Inches

_n = [0]

def tf(s,l,t,w,h,anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(l,t,w,h); f = tb.text_frame
    f.word_wrap = True
    f.margin_left=f.margin_right=f.margin_top=f.margin_bottom=0
    f.vertical_anchor = anchor
    return f

def run(p,text,size,color,bold=False,font=SANS,spc=None,italic=False):
    r = p.add_run(); r.text = text
    ft = r.font; ft.size=Pt(size); ft.bold=bold; ft.italic=italic
    ft.name=font; ft.color.rgb=color
    if spc is not None: r._r.get_or_add_rPr().set('spc', str(spc))
    return r

def para(f, first=False, align=PP_ALIGN.LEFT, before=0, after=0, line=None):
    p = f.paragraphs[0] if first else f.add_paragraph()
    p.alignment = align
    if before: p.space_before = Pt(before)
    if after:  p.space_after  = Pt(after)
    if line:   p.line_spacing = line
    return p

def box(s,l,t,w,h,fill,line=None,line_w=1.0,shadow=None,shape=MSO_SHAPE.RECTANGLE):
    if shadow is not None:
        sh = s.shapes.add_shape(shape, l+shadow, t+shadow, w, h)
        sh.fill.solid(); sh.fill.fore_color.rgb = SHADOW
        sh.line.fill.background(); sh.shadow.inherit=False
    r = s.shapes.add_shape(shape, l,t,w,h)
    if fill is None: r.fill.background()
    else: r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is None: r.line.fill.background()
    else: r.line.color.rgb = line; r.line.width = Pt(line_w)
    r.shadow.inherit = False
    return r

def hline(s,x,y,w,color=HAIR,pt=1.0):
    return box(s, I(x), I(y), I(w), Pt(pt).inches and I(pt/72.0), color)

def ticks(s, dark=False):
    """Registration tick marks in the four corners. Signature element."""
    c = GHOST_D if dark else HAIR
    m = 0.32; L = 0.16; t = 0.016
    for (cx,cy,hx,vy) in [(m,m,1,1),(13.333-m,m,-1,1),(m,7.5-m,1,-1),(13.333-m,7.5-m,-1,-1)]:
        box(s, I(min(cx,cx+hx*L)), I(cy-t/2), I(L), I(t), c)
        box(s, I(cx-t/2), I(min(cy,cy+vy*L)), I(t), I(L), c)

def ghost_num(s, dark=False, x=10.1, y=0.35):
    f = tf(s, I(x), I(y), I(3.0), I(2.6))
    run(para(f,first=True,align=PP_ALIGN.RIGHT), f"{_n[0]:02d}", 150, GHOST_D if dark else GHOST, bold=True, font=MONO)

def slide(dark=False, ghost=True, tick=True):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0,W,H)
    bg.fill.solid(); bg.fill.fore_color.rgb = INK if dark else PAPER
    bg.line.fill.background(); bg.shadow.inherit = False
    _n[0] += 1
    if ghost: ghost_num(s, dark)
    if tick: ticks(s, dark)
    return s

def eyebrow(s,l,t,text,color=MUTED,w=9.0):
    f = tf(s,l,t,I(w),I(0.28))
    run(para(f,first=True), text.upper(), 10.5, color, font=MONO, spc=240)

def dot(s,cx,cy,r,fill,line=None,lw=1.0):
    return box(s, I(cx-r), I(cy-r), I(2*r), I(2*r), fill, line=line, line_w=lw, shape=MSO_SHAPE.OVAL)

def pagefoot(s, dark=False):
    f = tf(s, I(0.55), I(7.08), I(6.0), I(0.3))
    run(para(f,first=True), "REVIVE · SEED", 8.5, GHOST_D if dark else FAINT, font=MONO, spc=200)

def case_card(s, x, y, w, detail=True, on_dark=False):
    h = w*0.47
    box(s, I(x), I(y), I(w), I(h), WHITE, line=None if on_dark else INK, line_w=1.5,
        shadow=None if on_dark else I(0.06))
    tb_h = h*0.115
    box(s, I(x), I(y), I(w), I(tb_h), INK2 if on_dark else INK)
    for i in range(3):
        dot(s, x+0.30+i*0.22, y+tb_h/2, 0.045, SLATE)
    f = tf(s, I(x+1.1), I(y), I(w-1.4), I(tb_h), anchor=MSO_ANCHOR.MIDDLE)
    run(para(f,first=True), "REVIVE · CONSOLE", 9 if not detail else 11, SLATE, font=MONO, spc=120)
    px = w*0.055
    if detail:
        f = tf(s, I(x+px), I(y+tb_h+0.28), I(w*0.5), I(0.45))
        run(para(f,first=True), "case_a1b2f9c3", 16, INK, bold=True, font=MONO)
        pw, ph = 1.9, 0.44
        box(s, I(x+w-px-pw), I(y+tb_h+0.26), I(pw), I(ph), GREEN_WASH, line=GREEN, line_w=1.25, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        f = tf(s, I(x+w-px-pw), I(y+tb_h+0.26), I(pw), I(ph), anchor=MSO_ANCHOR.MIDDLE)
        run(para(f,first=True,align=PP_ALIGN.CENTER), "RESUMED", 11, GREEN, bold=True, font=MONO, spc=100)
        f = tf(s, I(x+px), I(y+tb_h+0.85), I(w-2*px), I(0.4))
        p = para(f,first=True)
        run(p, "Action  ", 11, FAINT, font=MONO); run(p, "charge_customer", 13, INK, bold=True)
        run(p, "   ·   $49.00   ·   Stripe", 13, MUTED)
        ty = y+tb_h+2.0
    else:
        pw, ph = 1.25, 0.34
        box(s, I(x+w-px-pw), I(y+tb_h+0.16), I(pw), I(ph), GREEN_WASH, line=GREEN, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        f = tf(s, I(x+w-px-pw), I(y+tb_h+0.16), I(pw), I(ph), anchor=MSO_ANCHOR.MIDDLE)
        run(para(f,first=True,align=PP_ALIGN.CENTER), "RESUMED", 8.5, GREEN, bold=True, font=MONO, spc=80)
        f = tf(s, I(x+px), I(y+tb_h+0.18), I(w*0.5), I(0.35))
        run(para(f,first=True), "case_a1b2", 11, INK, bold=True, font=MONO)
        ty = y+tb_h+0.85
    states = ["Detected","Parked","Verified","Checked","Resumed"]
    span = w-2*px; step = span/(len(states)-1)
    box(s, I(x+px), I(ty-0.01), I(span), I(0.022), HAIR)
    for i,st in enumerate(states):
        cxp = x+px+i*step
        last = (i==len(states)-1)
        dot(s, cxp, ty, 0.075 if detail else 0.055, GREEN if last else COBALT)
        if detail:
            f = tf(s, I(cxp-0.7), I(ty+0.14), I(1.4), I(0.3))
            run(para(f,first=True,align=PP_ALIGN.CENTER), st, 9.5, INK if last else MUTED, font=MONO, bold=last)
    ly = ty + (0.75 if detail else 0.35)
    lh = 0.72 if detail else 0.45
    box(s, I(x+px), I(ly), I(w-2*px), I(lh), GREEN_WASH, line=GREEN, line_w=1.25 if detail else 1.0)
    f = tf(s, I(x+px+0.25), I(ly), I(w-2*px-0.4), I(lh), anchor=MSO_ANCHOR.MIDDLE)
    p = para(f,first=True)
    run(p, "✓  ", 15 if detail else 12, GREEN, bold=True)
    if detail:
        run(p, "Charge already went through. Revive did ", 13, INK, bold=True)
        run(p, "not", 13, GREEN, bold=True)
        run(p, " run it again.", 13, INK, bold=True)
    else:
        run(p, "not charged twice", 11, INK, bold=True)

# ================= 1 COVER : asymmetric split, ink panel bleeds right =================
s = slide(ghost=False, tick=False)
_n[0] = 1
# right ink panel, full-bleed to edges
box(s, I(8.35), 0, I(13.333-8.35), H, INK)
case_card(s, 8.85, 2.35, 3.9, detail=False, on_dark=True)
f = tf(s, I(8.85), I(5.35), I(3.9), I(1.2))
p = para(f,first=True,line=1.3)
run(p, "A real recovery record.\n", 12, SLATE)
run(p, "Live at revivelabs.app", 12, WHITE, bold=True)
# left content
ticks(s)
eyebrow(s, I(0.9), I(0.95), "Agent recovery control plane", COBALT_DEEP)
f = tf(s, I(0.82), I(1.55), I(7.4), I(2.0))
p = para(f,first=True)
run(p, "Revive", 108, INK, bold=True); run(p, ".", 108, COBALT)
f2 = tf(s, I(0.9), I(3.85), I(6.9), I(2.0))
p = para(f2,first=True,line=1.14)
run(p, "AI agents are starting to act:\nsending, buying, paying.", 22, INK, bold=True)
p2 = para(f2, line=1.2, before=12)
run(p2, "When one fails, Revive picks it back up ", 16, MUTED)
run(p2, "without doing it twice.", 16, COBALT_DEEP, bold=True)
f3 = tf(s, I(0.9), I(6.55), I(7.0), I(0.4))
p = para(f3,first=True)
run(p, "PRE-SEED", 10, FAINT, font=MONO, spc=220)
run(p, "     pip install revive-sdk     certified in production", 10.5, MUTED, font=MONO)

# ================= 2 WHY NOW : typographic scale ramp =================
s = slide()
eyebrow(s, I(0.9), I(0.85), "Why now", COBALT_DEEP)
f = tf(s, I(0.86), I(1.3), I(10.5), I(1.0))
run(para(f,first=True), "Software just started acting for us.", 34, INK, bold=True)
# ramp: three words, baseline aligned, growing
base = 4.35
words = [
    (0.95, "TALK", 30, MUTED,  "2023 to 2024", "Chatbots and copilots.", "A wrong answer is annoying."),
    (4.35, "ACT", 52, INK,    "2026", "Agents send, buy, pay.", "A wrong action is real damage."),
    (8.15, "SCALE", 74, COBALT_DEEP, "next", "Millions acting daily.", "Failure becomes a daily event."),
]
# rising connector
for i,(x,wd,sz,col,tag,l1,l2) in enumerate(words):
    f = tf(s, I(x), I(base-sz/48.0), I(4.9), I(sz/44.0+0.3), anchor=MSO_ANCHOR.BOTTOM)
    f.word_wrap = False
    run(para(f,first=True), wd, sz, col, bold=True, spc=40)
    box(s, I(x+0.02), I(base+0.12), I(3.0 if i==2 else 2.6), I(0.03), col if i==2 else HAIR)
    eyebrow(s, I(x+0.02), I(base+0.3), tag, col if i==2 else FAINT, w=3.2)
    f2 = tf(s, I(x+0.02), I(base+0.66), I(3.3), I(1.0))
    p = para(f2,first=True,line=1.15)
    run(p, l1+"\n", 13, INK, bold=True)
    run(p, l2, 12, MUTED)
box(s, I(0.9), I(6.3), I(11.5), I(0.78), INK, shadow=I(0.05))
f = tf(s, I(1.25), I(6.3), I(10.9), I(0.78), anchor=MSO_ANCHOR.MIDDLE)
p = para(f,first=True)
run(p, "Payments got Stripe. Outages got PagerDuty. ", 14.5, SLATE)
run(p, "Agents acting in the real world get Revive.", 14.5, WHITE, bold=True)
pagefoot(s)

# ================= 3 PROBLEM : full-dark terminal =================
s = slide(dark=True)
eyebrow(s, I(0.9), I(0.85), "The problem", RED_BR)
f = tf(s, I(0.86), I(1.3), I(11.6), I(1.7))
p = para(f,first=True,line=1.05)
run(p, "It's 3am. Your agent just\ncharged the customer ", 34, WHITE, bold=True)
run(p, "twice.", 34, RED_BR, bold=True)
# terminal log
tx, ty, tw, th = 0.9, 3.35, 7.6, 3.1
box(s, I(tx), I(ty), I(tw), I(th), INK2, line=GHOST_D, line_w=1.0)
logs = [
    ("03:07:41", "agent/invoice-runner   step 4 of 7", SLATE, False),
    ("03:07:41", "POST /v1/charges   $49.00 ............ ok", SLATE, False),
    ("03:07:44", "credential expired mid-run", RGBColor(0xd9,0xa4,0x3f), False),
    ("03:07:52", "retry #1   POST /v1/charges   $49.00 .. ok", WHITE, False),
    ("03:07:52", "DUPLICATE CHARGE. customer billed twice.", RED_BR, True),
]
ly = ty+0.35
for tstamp, msg, col, hl in logs:
    if hl:
        box(s, I(tx+0.18), I(ly-0.07), I(tw-0.36), I(0.42), RGBColor(0x33,0x1d,0x1e))
    f = tf(s, I(tx+0.35), I(ly), I(tw-0.6), I(0.35))
    p = para(f,first=True)
    run(p, tstamp+"  ", 11.5, RGBColor(0x55,0x5f,0x74), font=MONO)
    run(p, msg, 11.5, col, font=MONO, bold=hl)
    ly += 0.52
# right rail: three consequences
rx = 9.0
cons = [("IF YOU RETRY","it sends twice"),("IF YOU SKIP","the job breaks"),("IF YOU RECONNECT BLIND","wrong person acts")]
cy = 3.35
for lbl, txt in cons:
    box(s, I(rx), I(cy), I(0.05), I(0.82), RED_BR)
    f = tf(s, I(rx+0.25), I(cy), I(3.3), I(0.9))
    p = para(f,first=True)
    run(p, lbl+"\n", 9.5, RED_BR, font=MONO, spc=140)
    p2 = para(f, before=2)
    run(p2, txt, 15, WHITE, bold=True)
    cy += 1.12
pagefoot(s, dark=True)

# ================= 4 SOLUTION : rail with circle nodes =================
s = slide()
eyebrow(s, I(0.9), I(0.85), "The solution", COBALT_DEEP)
f = tf(s, I(0.86), I(1.3), I(11.0), I(1.0))
run(para(f,first=True), "A safety net under every action.", 34, INK, bold=True)
rail_y = 3.75
box(s, I(1.35), I(rail_y-0.012), I(10.6), I(0.024), HAIR)
steps = [("1","Spot","Know it failed"),("2","Hold","Freeze safely, lose nothing"),
         ("3","Verify","Right person reconnects"),("4","Check","Did it already happen?"),
         ("5","Continue","Picks up where it left off")]
xs = [1.35, 3.85, 6.35, 8.85, 11.35]
for i,((num,title,sub),cx) in enumerate(zip(steps,xs)):
    last = (i==4)
    r = 0.34
    dot(s, cx, rail_y, r+0.07, PAPER)
    dot(s, cx, rail_y, r, COBALT if last else PANEL, line=COBALT if last else INK, lw=1.5)
    f = tf(s, I(cx-r), I(rail_y-r), I(2*r), I(2*r), anchor=MSO_ANCHOR.MIDDLE)
    run(para(f,first=True,align=PP_ALIGN.CENTER), num, 17, WHITE if last else INK, bold=True, font=MONO)
    f2 = tf(s, I(cx-1.15), I(rail_y+0.55), I(2.3), I(1.3))
    p = para(f2,first=True,align=PP_ALIGN.CENTER)
    run(p, title, 16.5, COBALT_DEEP if last else INK, bold=True)
    p2 = para(f2, align=PP_ALIGN.CENTER, before=3, line=1.1)
    run(p2, sub, 11.5, MUTED)
f = tf(s, I(0.9), I(6.0), I(11.5), I(0.5))
p = para(f,first=True)
run(p, "ONE LINE OF CODE  ", 10.5, COBALT_DEEP, font=MONO, spc=180)
run(p, "Revive owns everything after the failure.", 14, INK, bold=True)
pagefoot(s)

# ================= 5 PRODUCT : console hero on ink band =================
s = slide(ghost=False)
box(s, 0, I(2.2), W, I(5.3), INK)
eyebrow(s, I(0.9), I(0.85), "Product", COBALT_DEEP)
f = tf(s, I(0.86), I(1.25), I(11.4), I(0.8))
run(para(f,first=True), "This is what a save looks like.", 32, INK, bold=True)
case_card(s, 1.97, 2.6, 9.4, detail=True, on_dark=True)

# ================= 6 WHY ONLY US : statement slide =================
s = slide()
eyebrow(s, I(0.9), I(0.95), "Why only us", COBALT_DEEP)
f = tf(s, I(0.86), I(1.7), I(11.8), I(2.6))
p = para(f,first=True,line=1.04)
run(p, "Anyone can retry.\n", 52, MUTED, bold=True)
run(p, "Only we know if it\nalready happened.", 52, INK, bold=True)
box(s, I(0.92), I(5.05), I(2.6), I(0.045), COBALT)
moat = [("We ask the source","We check the real system before anything runs again."),
        ("The right human","Same person proven back before the agent gets new keys."),
        ("We never hold keys","Logins stay in your vault. We store nothing.")]
x = 0.9
for title,sub in moat:
    f = tf(s, I(x), I(5.45), I(3.7), I(1.3))
    p = para(f,first=True)
    run(p, title, 15.5, INK, bold=True)
    p2 = para(f, before=4, line=1.15)
    run(p2, sub, 11.5, MUTED)
    x += 3.95
pagefoot(s)

# ================= 7 COMPETITION : bridge diagram =================
s = slide()
eyebrow(s, I(0.9), I(0.85), "Competition", COBALT_DEEP)
f = tf(s, I(0.86), I(1.3), I(11.6), I(1.6))
run(para(f,first=True,line=1.05), "Everyone owns one piece.\nNobody owns the recovery.", 32, INK, bold=True)
gy = 4.1; gh = 1.7
# two shore blocks
box(s, I(0.9), I(gy), I(4.1), I(gh), PANEL2, line=HAIR, line_w=1.0)
f = tf(s, I(0.9), I(gy+0.25), I(4.1), I(0.3), anchor=MSO_ANCHOR.TOP)
run(para(f,first=True,align=PP_ALIGN.CENTER), "RUNS THE AGENTS", 10, FAINT, font=MONO, spc=140)
f = tf(s, I(0.9), I(gy), I(4.1), I(gh), anchor=MSO_ANCHOR.MIDDLE)
run(para(f,first=True,align=PP_ALIGN.CENTER), "Temporal · LangGraph", 15, INK, bold=True)
box(s, I(8.33), I(gy), I(4.1), I(gh), PANEL2, line=HAIR, line_w=1.0)
f = tf(s, I(8.33), I(gy+0.25), I(4.1), I(0.3))
run(para(f,first=True,align=PP_ALIGN.CENTER), "LOGS THEM IN", 10, FAINT, font=MONO, spc=140)
f = tf(s, I(8.33), I(gy), I(4.1), I(gh), anchor=MSO_ANCHOR.MIDDLE)
run(para(f,first=True,align=PP_ALIGN.CENTER), "Nango · Auth0", 15, INK, bold=True)
# gap hatching
f = tf(s, I(5.0), I(gy+gh-0.42), I(3.33), I(0.35))
run(para(f,first=True,align=PP_ALIGN.CENTER), "· · · the gap · · ·", 11, FAINT, font=MONO)
# bridge block overlapping both
box(s, I(4.45), I(gy-0.75), I(4.45), I(1.35), INK, shadow=I(0.07))
f = tf(s, I(4.45), I(gy-0.75), I(4.45), I(1.35), anchor=MSO_ANCHOR.MIDDLE)
p = para(f,first=True,align=PP_ALIGN.CENTER)
run(p, "REVIVE", 11, LILAC, font=MONO, spc=220)
p2 = para(f, align=PP_ALIGN.CENTER, before=5)
run(p2, "the bridge between them", 15, WHITE, bold=True)
f = tf(s, I(0.9), I(6.3), I(11.5), I(0.5))
p = para(f,first=True)
run(p, "These two sides don't talk to each other. ", 13.5, INK, bold=True)
run(p, "Teams hand-roll half-fixes today. Revive is the real one.", 13.5, MUTED)
pagefoot(s)

# ================= 8 MARKET : nested squares =================
s = slide()
eyebrow(s, I(0.9), I(0.85), "Market", COBALT_DEEP)
f = tf(s, I(0.86), I(1.3), I(11.6), I(0.9))
p = para(f,first=True)
run(p, "Bottom-up, and it grows ", 32, INK, bold=True)
run(p, "tenfold", 32, COBALT_DEEP, bold=True)
run(p, " on its own.", 32, INK, bold=True)
# nested squares anchored bottom-left
bx, by = 0.95, 6.75   # common bottom-left corner
sq = [("TAM", 4.15, PANEL, HAIR), ("SAM", 2.7, COBALT_WASH, COBALT), ("SOM", 1.35, INK, INK)]
for tag, size, fill, line in sq:
    box(s, I(bx), I(by-size), I(size), I(size), fill, line=line, line_w=1.25)
f = tf(s, I(bx+0.18), I(by-4.15+0.12), I(1.6), I(0.35))
run(para(f,first=True), "TAM", 13, MUTED, bold=True, font=MONO, spc=80)
f = tf(s, I(bx+0.18), I(by-2.7+0.1), I(1.4), I(0.35))
run(para(f,first=True), "SAM", 13, COBALT_DEEP, bold=True, font=MONO, spc=80)
f = tf(s, I(bx), I(by-1.35), I(1.35), I(1.35), anchor=MSO_ANCHOR.MIDDLE)
run(para(f,first=True,align=PP_ALIGN.CENTER), "SOM", 13, WHITE, bold=True, font=MONO, spc=80)
# leader labels right
rows = [
    ("TAM", "$5B+", "Every team running production agents", "2M teams by 2030 × ~$2.5k/yr", by-4.15+0.35),
    ("SAM", "$100M → $1B+", "Agents touching email, money, CRM", "50k teams today × ~$2k/yr", by-2.7+0.5),
    ("SOM", "$5M ARR", "Our 3-year capture", "2.5k teams × ~$2k/yr", by-1.35+0.35),
]
for tag, val, lab, math, yy in rows:
    box(s, I(bx+4.35), I(yy+0.16), I(1.0), I(0.018), HAIR)
    f = tf(s, I(bx+5.55), I(yy-0.28), I(6.6), I(1.0))
    p = para(f,first=True)
    run(p, val+"   ", 24, COBALT_DEEP if tag!="TAM" else INK, bold=True)
    p2 = para(f, before=2)
    run(p2, lab+"  ", 13.5, INK, bold=True)
    run(p2, math, 11, FAINT, font=MONO)
f = tf(s, I(6.5), I(6.45), I(6.0), I(0.7))
p = para(f,first=True,line=1.2)
run(p, "Gartner puts agents in a third of enterprise software by 2028. ", 12, MUTED)
run(p, "Every new acting agent lands in our SAM.", 12, INK, bold=True)
pagefoot(s)

# ================= 9 BUSINESS MODEL : giant center tier =================
s = slide()
eyebrow(s, I(0.9), I(0.85), "Business model", COBALT_DEEP)
f = tf(s, I(0.86), I(1.3), I(11.4), I(0.8))
run(para(f,first=True), "We get paid every time we save one.", 32, INK, bold=True)
# center: TEAM hero card
cwx, cwy, cww, cwh = 4.72, 2.5, 3.9, 3.6
box(s, I(cwx), I(cwy), I(cww), I(cwh), INK, shadow=I(0.08))
eyebrow(s, I(cwx+0.4), I(cwy+0.4), "Team", LILAC, w=cww-0.7)
f = tf(s, I(cwx+0.35), I(cwy+0.85), I(cww-0.7), I(1.5))
p = para(f,first=True)
run(p, "$99", 66, WHITE, bold=True)
run(p, " /mo", 18, SLATE)
f = tf(s, I(cwx+0.4), I(cwy+2.35), I(cww-0.75), I(1.0))
p = para(f,first=True,line=1.25)
run(p, "The landing plan.\n", 14, WHITE, bold=True)
run(p, "25 connections · 10k saves a month", 12, SLATE)
# flanks
for (fx, name, price, note) in [(1.15,"DEV","$20 /mo","first workflow"),(9.35,"ENTERPRISE","Custom","where margin lives")]:
    box(s, I(fx), I(3.15), I(2.9), I(2.3), PANEL, line=HAIR, line_w=1.0)
    eyebrow(s, I(fx+0.3), I(3.45), name, MUTED, w=2.5)
    f = tf(s, I(fx+0.28), I(3.85), I(2.4), I(0.8))
    run(para(f,first=True), price, 26, INK, bold=True)
    f = tf(s, I(fx+0.3), I(4.7), I(2.4), I(0.5))
    run(para(f,first=True), note, 11.5, FAINT, font=MONO)
f = tf(s, I(0.9), I(6.5), I(11.5), I(0.5))
run(para(f,first=True), "More agents means more saves means a bigger plan. Upgrades happen on their own. Software margins.", 12.5, MUTED)
pagefoot(s)

# ================= 10 GTM : split channels | econ panel =================
s = slide()
eyebrow(s, I(0.9), I(0.85), "Go-to-market", COBALT_DEEP)
f = tf(s, I(0.86), I(1.3), I(11.4), I(0.9))
run(para(f,first=True), "Land where the builders already are.", 32, INK, bold=True)
# left: channel list
chan = [("01","Framework communities","devs hit this exact failure"),
        ("02","Vault partnerships","a connection dies, we step in"),
        ("03","Incident inbound","every agent mishap brings leads")]
cy = 2.75
for num,title,sub in chan:
    box(s, I(0.92), I(cy+0.06), I(0.045), I(0.95), COBALT)
    f = tf(s, I(1.25), I(cy), I(5.9), I(1.1))
    p = para(f,first=True)
    run(p, num+"   ", 13, COBALT_DEEP, bold=True, font=MONO)
    run(p, title, 17, INK, bold=True)
    p2 = para(f, before=4)
    run(p2, sub, 12.5, MUTED)
    cy += 1.28
# right: econ ink panel
ex, ey, ew, eh = 7.7, 2.6, 4.7, 3.55
box(s, I(ex), I(ey), I(ew), I(eh), INK, shadow=I(0.07))
eyebrow(s, I(ex+0.4), I(ey+0.35), "Unit economics · targets", LILAC, w=ew-0.7)
f = tf(s, I(ex+0.38), I(ey+0.75), I(ew-0.76), I(1.2))
p = para(f,first=True)
run(p, "5 : 1", 56, WHITE, bold=True)
run(p, "  LTV to CAC", 14, SLATE)
rows = [("CAC","< $500","community-led, not sales"),("LTV","~$2.4k","Team plan, 2-year retention")]
ry = ey+2.15
for k,v,note in rows:
    f = tf(s, I(ex+0.4), I(ry), I(ew-0.8), I(0.5))
    p = para(f,first=True)
    run(p, k+"  ", 11, LILAC, font=MONO, spc=100, bold=True)
    run(p, v+"   ", 16, WHITE, bold=True)
    run(p, note, 11, SLATE)
    ry += 0.6
f = tf(s, I(0.9), I(6.6), I(11.5), I(0.5))
p = para(f,first=True)
run(p, "12-MONTH TARGET  ", 10.5, COBALT_DEEP, font=MONO, spc=160)
run(p, "40 paying teams and a first enterprise logo, about $50k ARR. Pre-revenue today.", 12.5, MUTED)
pagefoot(s)

# ================= 11 TEAM =================
s = slide()
eyebrow(s, I(0.9), I(0.85), "Team", COBALT_DEEP)
f = tf(s, I(0.86), I(1.3), I(11.6), I(0.9))
run(para(f,first=True), "Three technical founders. All ship code.", 32, INK, bold=True)
team = [
    ("Srikanth Samy","SS","Co-founder & CEO","UC Berkeley","Built the entire system and proved it in production. Leads product and engineering."),
    ("Revanth Guda","RG","Co-founder & CTO","UCLA","Owns infrastructure and provider integrations, the connectors that make every platform recoverable."),
    ("Aarush Parekh","AP","Co-founder · Engineering","UC Santa Cruz","Owns SDKs and developer experience, the one-line install that gets Revive into every agent stack."),
]
cw, gap, x, top, ch = 3.66, 0.32, 0.9, 2.55, 3.7
PHOTO_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "photos")
for name,ini,title,school,blurb in team:
    box(s, I(x), I(top), I(cw), I(ch), PANEL, line=INK, line_w=1.25, shadow=I(0.055))
    # ghost initials watermark
    f = tf(s, I(x+cw-1.75), I(top+0.12), I(1.6), I(1.2))
    run(para(f,first=True,align=PP_ALIGN.RIGHT), ini, 58, PANEL2, bold=True, font=MONO)
    pr = 0.55
    photo = next((p2 for ext in ("png","jpg","jpeg") if os.path.exists(p2 := os.path.join(PHOTO_DIR, ini.lower()+"."+ext))), None)
    if photo:
        s.shapes.add_picture(photo, I(x+0.32), I(top+0.35), I(2*pr), I(2*pr))
    else:
        box(s, I(x+0.32), I(top+0.35), I(2*pr), I(2*pr), COBALT_WASH, line=COBALT, line_w=1.25, shape=MSO_SHAPE.OVAL)
        f = tf(s, I(x+0.32), I(top+0.35), I(2*pr), I(2*pr), anchor=MSO_ANCHOR.MIDDLE)
        run(para(f,first=True,align=PP_ALIGN.CENTER), ini, 22, COBALT_DEEP, bold=True, font=MONO)
    f = tf(s, I(x+0.3), I(top+1.7), I(cw-0.55), I(0.5))
    run(para(f,first=True,line=1.0), name, 19, INK, bold=True)
    eyebrow(s, I(x+0.32), I(top+2.2), title, COBALT_DEEP, w=cw-0.55)
    f2 = tf(s, I(x+0.3), I(top+2.55), I(cw-0.55), I(0.3))
    run(para(f2,first=True), school, 11, MUTED, font=MONO)
    f3 = tf(s, I(x+0.3), I(top+2.9), I(cw-0.55), I(0.75))
    run(para(f3,first=True,line=1.15), blurb, 11.5, MUTED)
    x += cw + gap
s.notes_slide.notes_text_frame.text = ("TO FINISH (not shown to investors): drop headshots into decks/photos/ss.png, rg.png, ap.png "
    "and rerun build_deck.py. Add one real credential per line for Revanth and Aarush. "
    "Confirm titles: Srikanth CEO, Revanth CTO, Aarush Engineering.")
pagefoot(s)

# ================= 12 VALIDATION : split done | next =================
s = slide(ghost=False)
# left ink panel
box(s, 0, 0, I(5.4), H, INK)
eyebrow(s, I(0.9), I(0.95), "Already done", LILAC)
done = ["Live product", "Works in production", "Developer tools shipped"]
dy = 1.7
for d in done:
    f = tf(s, I(0.9), I(dy), I(4.2), I(0.5))
    p = para(f,first=True)
    run(p, "✓  ", 16, GREEN_BR, bold=True)
    run(p, d, 16, WHITE, bold=True)
    dy += 0.62
f = tf(s, I(0.9), I(5.9), I(4.1), I(1.1))
p = para(f,first=True,line=1.3)
run(p, "Honest: no customers yet.\n", 13, WHITE, bold=True)
run(p, "This is the 60-day plan\nto change that.", 13, SLATE)
# right paper: plan
eyebrow(s, I(6.2), I(0.95), "Next 60 days", COBALT_DEEP)
f = tf(s, I(6.15), I(1.4), I(6.6), I(0.9))
run(para(f,first=True,line=1.05), "It's built. Now we prove\npeople want it.", 28, INK, bold=True)
plan = [("30+","interviews with agent teams","get the pain in their words"),
        ("5","design-partner pilots","free plan for a real case study"),
        ("","measure demand","installs · waitlist · community"),
        ("#1","publish the first pilot","the first real numbers")]
py = 3.15
for big,title,sub in plan:
    box(s, I(6.2), I(py+0.05), I(0.045), I(0.8), COBALT)
    f = tf(s, I(6.5), I(py), I(6.3), I(0.95))
    p = para(f,first=True)
    if big: run(p, big+"  ", 17, COBALT_DEEP, bold=True, font=MONO)
    run(p, title, 16, INK, bold=True)
    p2 = para(f, before=3)
    run(p2, sub, 12, MUTED)
    py += 1.02
pagefoot(s, dark=True)

# ================= 13 THE ASK : giant number =================
s = slide()
box(s, 0,0, I(0.16), H, COBALT)
eyebrow(s, I(0.9), I(0.9), "The ask", COBALT_DEEP)
f = tf(s, I(0.82), I(1.15), I(11.8), I(2.1))
p = para(f,first=True)
run(p, "$750k", 120, INK, bold=True)
f = tf(s, I(0.92), I(3.35), I(11.4), I(0.5))
p = para(f,first=True)
run(p, "18-MONTH RUNWAY", 13, COBALT_DEEP, font=MONO, spc=220, bold=True)
run(p, "      one risk left: distribution. Here is what the money turns it into.", 14, MUTED)
mile = [("5 → 10+","pilots turned into paying teams"),
        ("3 + 1","integrations live, first enterprise"),
        ("~$450k","ARR, ready to raise the A")]
x = 0.92
for big,lbl in mile:
    f = tf(s, I(x), I(4.25), I(3.7), I(1.3))
    p = para(f,first=True)
    run(p, big, 34, COBALT_DEEP, bold=True)
    p2 = para(f, before=5, line=1.1)
    run(p2, lbl, 13, INK, bold=True)
    box(s, I(x+0.02), I(5.55), I(3.3), I(0.035), HAIR)
    x += 3.85
box(s, I(0.92), I(6.1), I(11.45), I(0.85), INK, shadow=I(0.05))
f = tf(s, I(1.25), I(6.1), I(10.9), I(0.85), anchor=MSO_ANCHOR.MIDDLE)
p = para(f,first=True)
run(p, "THE MONEY BUYS   ", 10.5, LILAC, font=MONO, spc=160)
run(p, "runway for the team · first sales hire · two more integrations · security groundwork", 13, WHITE)
pagefoot(s)

# ================= 14 CLOSE =================
s = slide(dark=True, ghost=False, tick=True)
# ghost wordmark
f = tf(s, I(6.4), I(4.3), I(7.2), I(3.2))
p = para(f,first=True,align=PP_ALIGN.RIGHT)
run(p, "R.", 220, GHOST_D, bold=True)
eyebrow(s, I(0.9), I(1.5), "Revive Labs", LILAC)
f = tf(s, I(0.86), I(2.4), I(11.4), I(2.6))
p = para(f,first=True,line=1.02)
run(p, "Agents are getting hands.", 44, WHITE, bold=True)
p2 = para(f, line=1.02, before=6)
run(p2, "Someone has to own what happens\nwhen the hands slip.", 44, LILAC, bold=True)
f2 = tf(s, I(0.9), I(6.1), I(11.4), I(0.6))
p = para(f2,first=True)
run(p, "founders@revivelabs.app", 16, WHITE, bold=True)
run(p, "   ·   revivelabs.app   ·   deck and data room on request", 13.5, SLATE)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "revive-seed-deck.pptx")
prs.save(out)
print("saved", out, os.path.getsize(out), "bytes,", len(prs.slides._sldIdLst), "slides")

#!/usr/bin/env python3
"""Revive seed deck v9. HTML/CSS slides rendered via headless Chrome to PNG, assembled into pptx.
Fonts: Space Grotesk / Inter / JetBrains Mono (Google Fonts, fetched at render time)."""
import os, subprocess, glob

HERE = os.path.dirname(os.path.abspath(__file__))
SLIDES_DIR = os.path.join(HERE, "slides")
PNG_DIR = os.path.join(HERE, "png")
PHOTO_DIR = os.path.join(HERE, "photos")
os.makedirs(SLIDES_DIR, exist_ok=True)
os.makedirs(PNG_DIR, exist_ok=True)

FONTS_CSS = open(os.path.join(HERE, "fonts", "fonts-local.css")).read()

CSS = FONTS_CSS + """
*{margin:0;padding:0;box-sizing:border-box}
:root{
  --paper:#f4f5f1; --panel:#fbfcf8; --panel2:#eceee8;
  --ink:#12151d; --ink2:#1a1f2b; --ink3:#232a3a;
  --cobalt:#4967f2; --cobalt-deep:#2e49c8; --wash:#edf0ff;
  --muted:#66707e; --faint:#98a0aa; --hair:#d5d9d0;
  --red:#c2413a; --red-br:#e86a5f; --red-wash:#fcedeb;
  --green:#148060; --green-br:#3fbf8f; --green-wash:#e9f4ef;
  --lilac:#8ea0ff; --slate:#9aa6c8;
}
html,body{width:1920px;height:1080px;overflow:hidden}
body{font-family:'Inter',sans-serif;background:var(--paper);color:var(--ink);
  -webkit-font-smoothing:antialiased;position:relative}
body.dot{background-image:radial-gradient(circle,#dfe3da 1.4px,transparent 1.4px);background-size:34px 34px;background-position:-7px -7px}
body.dark{background:var(--ink);color:#fff}
body.dark.dot{background-image:radial-gradient(circle,#232a3a 1.4px,transparent 1.4px)}
.pad{position:absolute;inset:0;padding:110px 120px}
.eyebrow{font-family:'JetBrains Mono',monospace;font-size:15px;font-weight:500;
  letter-spacing:.42em;text-transform:uppercase;color:var(--cobalt-deep)}
.dark .eyebrow{color:var(--lilac)}
h1{font-family:'Space Grotesk',sans-serif;font-weight:700;letter-spacing:-.03em;line-height:1.02}
.disp{font-size:64px;margin-top:28px}
.mono{font-family:'JetBrains Mono',monospace}
.muted{color:var(--muted)} .faint{color:var(--faint)}
.tick{position:absolute;width:26px;height:26px;border-color:#c2c7bd !important}
.dark .tick{border-color:#2a3040 !important}
.tick.tl{top:52px;left:56px;border-top:2.5px solid;border-left:2.5px solid}
.tick.tr{top:52px;right:56px;border-top:2.5px solid;border-right:2.5px solid}
.tick.bl{bottom:52px;left:56px;border-bottom:2.5px solid;border-left:2.5px solid}
.tick.br{bottom:52px;right:56px;border-bottom:2.5px solid;border-right:2.5px solid}
.pageno{position:absolute;bottom:56px;right:110px;font-family:'JetBrains Mono',monospace;
  font-size:14px;letter-spacing:.3em;color:var(--faint)}
.dark .pageno{color:#3a4256}
.runner{position:absolute;bottom:56px;left:110px;font-family:'JetBrains Mono',monospace;
  font-size:13px;letter-spacing:.34em;color:#b9beb2}
.dark .runner{color:#2e3547}
.ghost{position:absolute;top:64px;right:96px;font-family:'Space Grotesk',sans-serif;font-weight:700;
  font-size:300px;line-height:1;color:#e7eae1;z-index:0;letter-spacing:-.04em}
.dark .ghost{color:#191e2b}
.z{position:relative;z-index:2}
.card{background:var(--panel);border:1.5px solid var(--ink);box-shadow:14px 14px 0 #e0e3da}
.dark .card{box-shadow:14px 14px 0 #0c0e14}
.chip{display:inline-flex;align-items:center;gap:14px;font-family:'JetBrains Mono',monospace;
  font-size:15px;letter-spacing:.18em}
.kicker{position:absolute;left:120px;right:120px;bottom:118px;background:var(--ink);color:#fff;
  padding:34px 44px;display:flex;align-items:center;gap:26px;box-shadow:14px 14px 0 #e0e3da}
.kicker .k{font-family:'JetBrains Mono',monospace;font-size:14px;letter-spacing:.34em;color:var(--lilac);white-space:nowrap}
"""

def console_card(mini=False, scale=1.0):
    if mini:
        return """
<div style="background:#fff;border-radius:14px;overflow:hidden;box-shadow:0 40px 80px -20px rgba(0,0,0,.55), 0 0 0 1px rgba(255,255,255,.08)">
  <div style="background:#1a1f2b;padding:16px 24px;display:flex;align-items:center;gap:9px">
    <span style="width:12px;height:12px;border-radius:50%;background:#e86a5f"></span>
    <span style="width:12px;height:12px;border-radius:50%;background:#e0b34c"></span>
    <span style="width:12px;height:12px;border-radius:50%;background:#3fbf8f"></span>
    <span class="mono" style="font-size:13px;letter-spacing:.28em;color:#9aa6c8;margin-left:14px">REVIVE · CONSOLE</span>
  </div>
  <div style="padding:28px 32px 30px;color:#12151d">
    <div style="display:flex;justify-content:space-between;align-items:center">
      <span class="mono" style="font-size:19px;font-weight:700">case_a1b2</span>
      <span class="mono" style="font-size:12px;font-weight:700;letter-spacing:.2em;color:#148060;background:#e9f4ef;border:1.5px solid #148060;border-radius:99px;padding:7px 16px">RESUMED</span>
    </div>
    <div style="margin-top:30px;position:relative;height:14px">
      <div style="position:absolute;top:6px;left:0;right:0;height:2.5px;background:#e2e5dd"></div>
      <div style="position:absolute;top:6px;left:0;width:100%;display:flex;justify-content:space-between">
        <span style="width:13px;height:13px;border-radius:50%;background:#4967f2;margin-top:-5px"></span>
        <span style="width:13px;height:13px;border-radius:50%;background:#4967f2;margin-top:-5px"></span>
        <span style="width:13px;height:13px;border-radius:50%;background:#4967f2;margin-top:-5px"></span>
        <span style="width:13px;height:13px;border-radius:50%;background:#4967f2;margin-top:-5px"></span>
        <span style="width:15px;height:15px;border-radius:50%;background:#148060;margin-top:-6px;box-shadow:0 0 0 5px #e9f4ef"></span>
      </div>
    </div>
    <div style="margin-top:26px;background:#e9f4ef;border:1.5px solid #148060;border-radius:10px;padding:14px 20px;font-size:15.5px;font-weight:600">
      <span style="color:#148060">✓</span>&nbsp; not charged twice
    </div>
  </div>
</div>"""
    return """
<div style="background:#fff;border-radius:18px;overflow:hidden;box-shadow:0 60px 120px -30px rgba(0,0,0,.6), 0 0 0 1px rgba(255,255,255,.06)">
  <div style="background:#1a1f2b;padding:22px 34px;display:flex;align-items:center;gap:11px">
    <span style="width:15px;height:15px;border-radius:50%;background:#e86a5f"></span>
    <span style="width:15px;height:15px;border-radius:50%;background:#e0b34c"></span>
    <span style="width:15px;height:15px;border-radius:50%;background:#3fbf8f"></span>
    <span class="mono" style="font-size:15px;letter-spacing:.3em;color:#9aa6c8;margin-left:18px">REVIVE · CONSOLE</span>
    <span class="mono" style="font-size:14px;color:#5b6478;margin-left:auto">revivelabs.app/app/runs</span>
  </div>
  <div style="padding:44px 52px 48px;color:#12151d">
    <div style="display:flex;justify-content:space-between;align-items:flex-start">
      <div>
        <div class="mono" style="font-size:27px;font-weight:700">case_a1b2f9c3</div>
        <div style="margin-top:14px;font-size:19px;color:#66707e">
          <span class="mono" style="font-size:15px;color:#98a0aa">ACTION</span>&nbsp;&nbsp;
          <b style="color:#12151d">charge_customer</b> &nbsp;·&nbsp; $49.00 &nbsp;·&nbsp; Stripe
        </div>
      </div>
      <span class="mono" style="font-size:15px;font-weight:700;letter-spacing:.22em;color:#148060;background:#e9f4ef;border:2px solid #148060;border-radius:99px;padding:12px 26px">RESUMED</span>
    </div>
    <div style="margin-top:52px;position:relative">
      <div style="position:absolute;top:8px;left:0;right:0;height:3px;background:#e2e5dd"></div>
      <div style="display:flex;justify-content:space-between;position:relative">
        <div style="text-align:center"><span style="display:block;width:18px;height:18px;border-radius:50%;background:#4967f2;margin:0 auto"></span><span class="mono" style="font-size:14px;color:#66707e;display:block;margin-top:14px">Detected</span></div>
        <div style="text-align:center"><span style="display:block;width:18px;height:18px;border-radius:50%;background:#4967f2;margin:0 auto"></span><span class="mono" style="font-size:14px;color:#66707e;display:block;margin-top:14px">Parked</span></div>
        <div style="text-align:center"><span style="display:block;width:18px;height:18px;border-radius:50%;background:#4967f2;margin:0 auto"></span><span class="mono" style="font-size:14px;color:#66707e;display:block;margin-top:14px">Verified</span></div>
        <div style="text-align:center"><span style="display:block;width:18px;height:18px;border-radius:50%;background:#4967f2;margin:0 auto"></span><span class="mono" style="font-size:14px;color:#66707e;display:block;margin-top:14px">Checked</span></div>
        <div style="text-align:center"><span style="display:block;width:22px;height:22px;border-radius:50%;background:#148060;margin:-2px auto 0;box-shadow:0 0 0 7px #e9f4ef"></span><span class="mono" style="font-size:14px;color:#12151d;font-weight:700;display:block;margin-top:12px">Resumed</span></div>
      </div>
    </div>
    <div style="margin-top:48px;background:#e9f4ef;border:2px solid #148060;border-radius:14px;padding:24px 32px;font-size:21px;font-weight:600">
      <span style="color:#148060;font-size:24px">✓</span>&nbsp;&nbsp;Charge already went through. Revive did <span style="color:#148060">not</span> run it again.
    </div>
  </div>
</div>"""

def frame(body, n, dark=False, dot=True, ghost=True, runner=True, ticks=True):
    cls = ("dark " if dark else "") + ("dot" if dot else "")
    g = f'<div class="ghost">{n:02d}</div>' if ghost else ""
    t = '<div class="tick tl"></div><div class="tick tr"></div><div class="tick bl"></div><div class="tick br"></div>' if ticks else ""
    r = '<div class="runner">REVIVE — SEED</div>' if runner else ""
    pn = f'<div class="pageno">{n:02d} / 14</div>' if n else ""
    return f"""<!doctype html><html><head><meta charset="utf-8"><style>{CSS}</style></head>
<body class="{cls}">{g}{t}{r}{pn}{body}</body></html>"""

S = []

# ---------- 1 COVER ----------
S.append(frame(f"""
<div style="position:absolute;top:0;right:0;bottom:0;width:645px;background:linear-gradient(160deg,#1a1f2b 0%,#12151d 60%);box-shadow:inset 1px 0 0 #232a3a">
  <div style="position:absolute;inset:0;background-image:radial-gradient(circle,#222939 1.4px,transparent 1.4px);background-size:34px 34px"></div>
  <div style="position:absolute;top:290px;left:70px;right:70px">{console_card(mini=True)}
    <div style="margin-top:44px;font-size:17px;color:#9aa6c8;line-height:1.6">A real recovery record.<br><b style="color:#fff">Live at revivelabs.app</b></div>
  </div>
</div>
<div class="pad">
  <div class="eyebrow">Agent recovery control plane</div>
  <h1 style="font-size:190px;margin-top:110px;letter-spacing:-.045em">Revive<span style="color:var(--cobalt)">.</span></h1>
  <div style="margin-top:70px;max-width:980px">
    <div style="font-family:'Space Grotesk';font-weight:700;font-size:46px;line-height:1.15;letter-spacing:-.02em">AI agents are starting to act:<br>sending, buying, paying.</div>
    <div style="margin-top:30px;font-size:27px;line-height:1.5;color:var(--muted)">When one fails, Revive picks it back up <b style="color:var(--cobalt-deep)">without doing it twice.</b></div>
  </div>
  <div style="position:absolute;bottom:112px;left:120px;display:flex;gap:52px" class="mono">
    <span style="font-size:15px;letter-spacing:.3em;color:var(--faint)">PRE-SEED</span>
    <span style="font-size:15px;letter-spacing:.12em;color:var(--muted)">pip install revive-sdk</span>
    <span style="font-size:15px;letter-spacing:.12em;color:var(--muted)">certified in production</span>
  </div>
</div>""", 1, ghost=False, runner=False))

# ---------- 2 WHY NOW ----------
S.append(frame(f"""
<div class="pad z">
  <div class="eyebrow">Why now</div>
  <h1 class="disp">Software just started acting for us.</h1>
  <div style="display:flex;align-items:flex-end;gap:70px;margin-top:150px;height:330px">
    <div style="flex:1">
      <div style="font-family:'Space Grotesk';font-weight:700;font-size:56px;color:#a7aeb8;letter-spacing:.02em">TALK</div>
      <div style="height:4px;background:var(--hair);margin:18px 0 22px"></div>
      <div class="mono" style="font-size:14px;letter-spacing:.3em;color:var(--faint)">2023 – 2024</div>
      <div style="font-size:22px;font-weight:600;margin-top:16px">Chatbots and copilots.</div>
      <div style="font-size:19px;color:var(--muted);margin-top:6px">A wrong answer is annoying.</div>
    </div>
    <div style="flex:1">
      <div style="font-family:'Space Grotesk';font-weight:700;font-size:96px;letter-spacing:.01em">ACT</div>
      <div style="height:4px;background:var(--ink);margin:18px 0 22px"></div>
      <div class="mono" style="font-size:14px;letter-spacing:.3em;color:var(--faint)">2026</div>
      <div style="font-size:22px;font-weight:600;margin-top:16px">Agents send, buy, pay.</div>
      <div style="font-size:19px;color:var(--muted);margin-top:6px">A wrong action is real damage.</div>
    </div>
    <div style="flex:1.25">
      <div style="font-family:'Space Grotesk';font-weight:700;font-size:150px;color:var(--cobalt);letter-spacing:-.01em;line-height:.9">SCALE</div>
      <div style="height:5px;background:var(--cobalt);margin:18px 0 22px"></div>
      <div class="mono" style="font-size:14px;letter-spacing:.3em;color:var(--cobalt-deep)">NEXT</div>
      <div style="font-size:22px;font-weight:600;margin-top:16px">Millions acting daily.</div>
      <div style="font-size:19px;color:var(--muted);margin-top:6px">Failure becomes a daily event.</div>
    </div>
  </div>
</div>
<div class="kicker"><span class="k">The pattern</span>
  <span style="font-size:24px;color:var(--slate)">Payments got Stripe. Outages got PagerDuty. <b style="color:#fff">Agents acting in the real world get Revive.</b></span>
</div>""", 2))

# ---------- 3 PROBLEM (terminal) ----------
S.append(frame(f"""
<div class="pad z">
  <div class="eyebrow" style="color:var(--red-br)">The problem</div>
  <h1 class="disp" style="color:#fff">It's 3am. Your agent just<br>charged the customer <span style="color:var(--red-br)">twice.</span></h1>
  <div style="display:flex;gap:80px;margin-top:96px">
    <div style="flex:1.6;background:#171c27;border:1px solid #262d3d;border-radius:16px;padding:46px 50px;box-shadow:0 50px 100px -40px rgba(0,0,0,.8)" class="mono">
      <div style="display:flex;gap:10px;margin-bottom:34px">
        <span style="width:13px;height:13px;border-radius:50%;background:#e86a5f"></span>
        <span style="width:13px;height:13px;border-radius:50%;background:#e0b34c"></span>
        <span style="width:13px;height:13px;border-radius:50%;background:#3fbf8f"></span>
      </div>
      <div style="font-size:20px;line-height:2.35;white-space:nowrap">
        <span style="color:#4a5468">03:07:41</span>&nbsp;&nbsp;<span style="color:#9aa6c8">agent/invoice-runner &nbsp;step 4 of 7</span><br>
        <span style="color:#4a5468">03:07:41</span>&nbsp;&nbsp;<span style="color:#9aa6c8">POST /v1/charges &nbsp;$49.00 ......... ok</span><br>
        <span style="color:#4a5468">03:07:44</span>&nbsp;&nbsp;<span style="color:#e0b34c">credential expired mid-run</span><br>
        <span style="color:#4a5468">03:07:52</span>&nbsp;&nbsp;<span style="color:#e6e9f0">retry #1 &nbsp;POST /v1/charges &nbsp;$49.00 .. ok</span><br>
        <span style="display:inline-block;background:#3a1f20;border-left:4px solid #e86a5f;padding:2px 18px;margin-left:-18px"><span style="color:#4a5468">03:07:52</span>&nbsp;&nbsp;<b style="color:#e86a5f">DUPLICATE CHARGE — customer billed twice</b></span>
      </div>
    </div>
    <div style="flex:1;display:flex;flex-direction:column;justify-content:center;gap:52px">
      <div style="border-left:5px solid var(--red-br);padding-left:28px">
        <div class="mono" style="font-size:14px;letter-spacing:.3em;color:var(--red-br)">IF YOU RETRY</div>
        <div style="font-family:'Space Grotesk';font-weight:700;font-size:33px;margin-top:10px">it sends twice</div>
      </div>
      <div style="border-left:5px solid var(--red-br);padding-left:28px">
        <div class="mono" style="font-size:14px;letter-spacing:.3em;color:var(--red-br)">IF YOU SKIP</div>
        <div style="font-family:'Space Grotesk';font-weight:700;font-size:33px;margin-top:10px">the job breaks</div>
      </div>
      <div style="border-left:5px solid var(--red-br);padding-left:28px">
        <div class="mono" style="font-size:14px;letter-spacing:.3em;color:var(--red-br)">IF YOU RECONNECT BLIND</div>
        <div style="font-family:'Space Grotesk';font-weight:700;font-size:33px;margin-top:10px">wrong person acts</div>
      </div>
    </div>
  </div>
</div>""", 3, dark=True))

# ---------- 4 SOLUTION ----------
steps_html = ""
labels = [("1","Spot","Know it failed"),("2","Hold","Freeze safely, lose nothing"),("3","Verify","Right person reconnects"),("4","Check","Did it already happen?"),("5","Continue","Picks up where it left off")]
for i,(n,t,sub) in enumerate(labels):
    last = i==4
    node = f"""<div style="width:96px;height:96px;border-radius:50%;display:flex;align-items:center;justify-content:center;
      {'background:var(--cobalt);color:#fff;box-shadow:0 20px 40px -12px rgba(73,103,242,.55)' if last else 'background:#fff;border:2.5px solid var(--ink)'};
      font-family:'Space Grotesk';font-weight:700;font-size:38px;margin:0 auto">{n}</div>"""
    steps_html += f"""<div style="flex:1;text-align:center;position:relative;z-index:2">{node}
      <div style="font-family:'Space Grotesk';font-weight:700;font-size:31px;margin-top:34px;{'color:var(--cobalt-deep)' if last else ''}">{t}</div>
      <div style="font-size:18.5px;color:var(--muted);margin-top:10px;padding:0 20px">{sub}</div></div>"""
S.append(frame(f"""
<div class="pad z">
  <div class="eyebrow">The solution</div>
  <h1 class="disp">A safety net under every action.</h1>
  <div style="position:relative;margin-top:170px">
    <div style="position:absolute;top:48px;left:9%;right:9%;height:3px;background:var(--hair)"></div>
    <div style="display:flex">{steps_html}</div>
  </div>
  <div style="margin-top:130px;display:flex;align-items:center;gap:30px">
    <span class="mono" style="font-size:15px;letter-spacing:.34em;color:var(--cobalt-deep);font-weight:700">ONE LINE OF CODE</span>
    <span style="font-size:24px;font-weight:600">Revive owns everything after the failure.</span>
  </div>
</div>""", 4))

# ---------- 5 PRODUCT ----------
S.append(frame(f"""
<div style="position:absolute;left:0;right:0;top:315px;bottom:0;background:linear-gradient(180deg,#1a1f2b,#12151d)"></div>
<div class="pad z">
  <div class="eyebrow">Product</div>
  <h1 class="disp">This is what a save looks like.</h1>
  <div style="max-width:1340px;margin:105px auto 0">{console_card()}</div>
</div>""", 5, ghost=False, runner=False))

# ---------- 6 WHY ONLY US ----------
S.append(frame(f"""
<div class="pad z">
  <div class="eyebrow">Why only us</div>
  <div style="font-family:'Space Grotesk';font-weight:700;font-size:104px;line-height:1.04;letter-spacing:-.03em;margin-top:70px">
    <span style="color:#a7aeb8">Anyone can retry.</span><br>Only we know if it<br>already happened.
  </div>
  <div style="width:250px;height:7px;background:var(--cobalt);margin-top:56px"></div>
  <div style="display:flex;gap:76px;margin-top:64px">
    <div style="flex:1"><div style="font-size:25px;font-weight:700">We ask the source</div>
      <div style="font-size:19px;color:var(--muted);margin-top:12px;line-height:1.55">We check the real system before anything runs again.</div></div>
    <div style="flex:1"><div style="font-size:25px;font-weight:700">The right human</div>
      <div style="font-size:19px;color:var(--muted);margin-top:12px;line-height:1.55">Same person proven back before the agent gets new keys.</div></div>
    <div style="flex:1"><div style="font-size:25px;font-weight:700">We never hold keys</div>
      <div style="font-size:19px;color:var(--muted);margin-top:12px;line-height:1.55">Logins stay in your vault. We store nothing.</div></div>
  </div>
</div>""", 6))

# ---------- 7 COMPETITION ----------
S.append(frame(f"""
<div class="pad z">
  <div class="eyebrow">Competition</div>
  <h1 class="disp">Everyone owns one piece.<br>Nobody owns the recovery.</h1>
  <div style="position:relative;margin-top:170px;height:360px">
    <div style="position:absolute;left:0;top:110px;width:560px;height:230px;background:var(--panel2);border:1.5px solid var(--hair);display:flex;flex-direction:column;align-items:center;justify-content:center;gap:16px">
      <span class="mono" style="font-size:14px;letter-spacing:.32em;color:var(--faint)">RUNS THE AGENTS</span>
      <span style="font-family:'Space Grotesk';font-weight:700;font-size:31px">Temporal · LangGraph</span>
    </div>
    <div style="position:absolute;right:0;top:110px;width:560px;height:230px;background:var(--panel2);border:1.5px solid var(--hair);display:flex;flex-direction:column;align-items:center;justify-content:center;gap:16px">
      <span class="mono" style="font-size:14px;letter-spacing:.32em;color:var(--faint)">LOGS THEM IN</span>
      <span style="font-family:'Space Grotesk';font-weight:700;font-size:31px">Nango · Auth0</span>
    </div>
    <div class="mono" style="position:absolute;left:0;right:0;top:305px;text-align:center;font-size:16px;letter-spacing:.3em;color:var(--faint)">· · · the gap · · ·</div>
    <div style="position:absolute;left:50%;transform:translateX(-50%);top:0;width:640px;height:190px;background:var(--ink);box-shadow:22px 22px 0 #e0e3da;display:flex;flex-direction:column;align-items:center;justify-content:center;gap:14px">
      <span class="mono" style="font-size:15px;letter-spacing:.42em;color:var(--lilac)">REVIVE</span>
      <span style="font-family:'Space Grotesk';font-weight:700;font-size:33px;color:#fff">the bridge between them</span>
    </div>
  </div>
  <div style="margin-top:88px;font-size:23px"><b>These two sides don't talk to each other.</b> <span class="muted">Teams hand-roll half-fixes today. Revive is the real one.</span></div>
</div>""", 7))

# ---------- 8 MARKET ----------
S.append(frame(f"""
<div class="pad z">
  <div class="eyebrow">Market</div>
  <h1 class="disp">Bottom-up, and it grows <span style="color:var(--cobalt)">tenfold</span> on its own.</h1>
  <div style="display:flex;gap:110px;margin-top:96px;align-items:flex-end">
    <div style="position:relative;width:600px;height:600px;flex:none">
      <div style="position:absolute;left:0;bottom:0;width:600px;height:600px;background:var(--panel);border:1.5px solid var(--hair)"></div>
      <div style="position:absolute;left:0;bottom:0;width:390px;height:390px;background:var(--wash);border:2px solid var(--cobalt)"></div>
      <div style="position:absolute;left:0;bottom:0;width:195px;height:195px;background:var(--ink);display:flex;align-items:center;justify-content:center">
        <span class="mono" style="font-size:19px;letter-spacing:.24em;color:#fff;font-weight:700">SOM</span></div>
      <span class="mono" style="position:absolute;left:26px;top:22px;font-size:19px;letter-spacing:.24em;color:var(--muted);font-weight:700">TAM</span>
      <span class="mono" style="position:absolute;left:26px;bottom:400px;font-size:19px;letter-spacing:.24em;color:var(--cobalt-deep);font-weight:700;margin-bottom:14px">SAM</span>
    </div>
    <div style="flex:1;display:flex;flex-direction:column;gap:66px;padding-bottom:12px">
      <div><div style="font-family:'Space Grotesk';font-weight:700;font-size:56px">$5B+</div>
        <div style="font-size:22px;font-weight:600;margin-top:6px">Every team running production agents <span class="mono faint" style="font-size:16px;font-weight:400">&nbsp;2M teams by 2030 × ~$2.5k/yr</span></div></div>
      <div><div style="font-family:'Space Grotesk';font-weight:700;font-size:56px;color:var(--cobalt-deep)">$100M → $1B+</div>
        <div style="font-size:22px;font-weight:600;margin-top:6px">Agents touching email, money, CRM <span class="mono faint" style="font-size:16px;font-weight:400">&nbsp;50k teams today × ~$2k/yr</span></div></div>
      <div><div style="font-family:'Space Grotesk';font-weight:700;font-size:56px;color:var(--cobalt-deep)">$5M ARR</div>
        <div style="font-size:22px;font-weight:600;margin-top:6px">Our 3-year capture <span class="mono faint" style="font-size:16px;font-weight:400">&nbsp;2.5k teams × ~$2k/yr</span></div></div>
      <div style="font-size:19px;color:var(--muted);line-height:1.5">Gartner puts agents in a third of enterprise software by 2028.<br><b style="color:var(--ink)">Every new acting agent lands in our SAM.</b></div>
    </div>
  </div>
</div>""", 8))

# ---------- 9 BUSINESS MODEL ----------
S.append(frame(f"""
<div class="pad z">
  <div class="eyebrow">Business model</div>
  <h1 class="disp">We get paid every time we save one.</h1>
  <div style="display:flex;align-items:center;justify-content:center;gap:56px;margin-top:110px">
    <div style="width:400px;background:var(--panel);border:1.5px solid var(--hair);padding:46px 46px 50px">
      <div class="mono" style="font-size:14px;letter-spacing:.32em;color:var(--muted)">DEV</div>
      <div style="font-family:'Space Grotesk';font-weight:700;font-size:64px;margin-top:22px">$20<span style="font-size:26px;color:var(--muted);font-weight:500"> /mo</span></div>
      <div class="mono" style="font-size:16px;color:var(--faint);margin-top:26px">first workflow</div>
    </div>
    <div style="width:520px;background:linear-gradient(165deg,#1c2230,#12151d);padding:60px 56px 64px;box-shadow:0 60px 110px -35px rgba(18,21,29,.75);border-radius:4px">
      <div class="mono" style="font-size:15px;letter-spacing:.36em;color:var(--lilac)">TEAM</div>
      <div style="font-family:'Space Grotesk';font-weight:700;font-size:130px;color:#fff;margin-top:16px;letter-spacing:-.03em">$99<span style="font-size:34px;color:var(--slate);font-weight:500;letter-spacing:0"> /mo</span></div>
      <div style="font-size:24px;color:#fff;font-weight:700;margin-top:26px">The landing plan.</div>
      <div style="font-size:19px;color:var(--slate);margin-top:10px">25 connections · 10k saves a month</div>
    </div>
    <div style="width:400px;background:var(--panel);border:1.5px solid var(--hair);padding:46px 46px 50px">
      <div class="mono" style="font-size:14px;letter-spacing:.32em;color:var(--muted)">ENTERPRISE</div>
      <div style="font-family:'Space Grotesk';font-weight:700;font-size:64px;margin-top:22px">Custom</div>
      <div class="mono" style="font-size:16px;color:var(--faint);margin-top:26px">where margin lives</div>
    </div>
  </div>
  <div style="text-align:center;margin-top:84px;font-size:21px;color:var(--muted)">More agents means more saves means a bigger plan. Upgrades happen on their own. Software margins.</div>
</div>""", 9))

# ---------- 10 GTM ----------
S.append(frame(f"""
<div class="pad z">
  <div class="eyebrow">Go-to-market</div>
  <h1 class="disp">Land where the builders already are.</h1>
  <div style="display:flex;gap:110px;margin-top:110px">
    <div style="flex:1.1;display:flex;flex-direction:column;gap:62px">
      <div style="border-left:6px solid var(--cobalt);padding-left:34px">
        <div><span class="mono" style="font-size:19px;color:var(--cobalt-deep);font-weight:700">01</span>
        <span style="font-family:'Space Grotesk';font-weight:700;font-size:33px;margin-left:18px">Framework communities</span></div>
        <div style="font-size:20px;color:var(--muted);margin-top:10px">devs hit this exact failure</div>
      </div>
      <div style="border-left:6px solid var(--cobalt);padding-left:34px">
        <div><span class="mono" style="font-size:19px;color:var(--cobalt-deep);font-weight:700">02</span>
        <span style="font-family:'Space Grotesk';font-weight:700;font-size:33px;margin-left:18px">Vault partnerships</span></div>
        <div style="font-size:20px;color:var(--muted);margin-top:10px">a connection dies, we step in</div>
      </div>
      <div style="border-left:6px solid var(--cobalt);padding-left:34px">
        <div><span class="mono" style="font-size:19px;color:var(--cobalt-deep);font-weight:700">03</span>
        <span style="font-family:'Space Grotesk';font-weight:700;font-size:33px;margin-left:18px">Incident inbound</span></div>
        <div style="font-size:20px;color:var(--muted);margin-top:10px">every agent mishap brings leads</div>
      </div>
    </div>
    <div style="flex:1;background:linear-gradient(165deg,#1c2230,#12151d);padding:58px 60px;box-shadow:0 60px 110px -35px rgba(18,21,29,.7);border-radius:4px;align-self:flex-start">
      <div class="mono" style="font-size:14px;letter-spacing:.34em;color:var(--lilac)">UNIT ECONOMICS · TARGETS</div>
      <div style="font-family:'Space Grotesk';font-weight:700;font-size:120px;color:#fff;margin-top:20px;letter-spacing:-.02em">5:1<span style="font-size:30px;color:var(--slate);font-weight:500;letter-spacing:0;margin-left:20px">LTV to CAC</span></div>
      <div style="margin-top:44px;display:flex;flex-direction:column;gap:24px;font-size:22px">
        <div><span class="mono" style="color:var(--lilac);font-size:16px;letter-spacing:.2em">CAC&nbsp;&nbsp;</span><b style="color:#fff">&lt; $500</b> <span style="color:var(--slate)">&nbsp;community-led, not sales</span></div>
        <div><span class="mono" style="color:var(--lilac);font-size:16px;letter-spacing:.2em">LTV&nbsp;&nbsp;</span><b style="color:#fff">~$2.4k</b> <span style="color:var(--slate)">&nbsp;Team plan, 2-year retention</span></div>
      </div>
    </div>
  </div>
  <div style="margin-top:96px;display:flex;align-items:center;gap:28px">
    <span class="mono" style="font-size:15px;letter-spacing:.32em;color:var(--cobalt-deep);font-weight:700">12-MONTH TARGET</span>
    <span style="font-size:22px;color:var(--muted)">40 paying teams and a first enterprise logo, about <b style="color:var(--ink)">$50k ARR</b>. Pre-revenue today.</span>
  </div>
</div>""", 10))

# ---------- 11 TEAM ----------
def founder(ini, name, title, school, blurb):
    photo = None
    for ext in ("png","jpg","jpeg"):
        p = os.path.join(PHOTO_DIR, ini.lower()+"."+ext)
        if os.path.exists(p): photo = p; break
    if photo:
        avatar = f'<img src="file://{photo}" style="width:150px;height:150px;border-radius:50%;object-fit:cover;border:3px solid var(--cobalt)">'
    else:
        avatar = f'''<div style="width:150px;height:150px;border-radius:50%;background:var(--wash);border:3px solid var(--cobalt);
          display:flex;align-items:center;justify-content:center;font-family:'Space Grotesk';font-weight:700;font-size:46px;color:var(--cobalt-deep)">{ini}</div>'''
    return f"""
    <div style="flex:1;background:var(--panel);border:1.5px solid var(--ink);box-shadow:16px 16px 0 #e0e3da;padding:52px 48px;position:relative;overflow:hidden">
      <div style="position:absolute;top:14px;right:26px;font-family:'Space Grotesk';font-weight:700;font-size:120px;color:#eef0e9;line-height:1">{ini}</div>
      <div style="position:relative">{avatar}
      <div style="font-family:'Space Grotesk';font-weight:700;font-size:37px;margin-top:38px">{name}</div>
      <div class="mono" style="font-size:15px;letter-spacing:.28em;color:var(--cobalt-deep);font-weight:700;margin-top:14px">{title}</div>
      <div class="mono" style="font-size:16px;color:var(--muted);margin-top:12px">{school}</div>
      <div style="font-size:19px;color:var(--muted);line-height:1.55;margin-top:22px">{blurb}</div></div>
    </div>"""
S.append(frame(f"""
<div class="pad z">
  <div class="eyebrow">Team</div>
  <h1 class="disp">Three technical founders. All ship code.</h1>
  <div style="display:flex;gap:54px;margin-top:96px">
    {founder("SS","Srikanth Samy","CO-FOUNDER & CEO","UC Berkeley","Built the entire system and proved it in production. Leads product and engineering.")}
    {founder("RG","Revanth Guda","CO-FOUNDER & CTO","UCLA","Owns infrastructure and provider integrations, the connectors that make every platform recoverable.")}
    {founder("AP","Aarush Parekh","CO-FOUNDER · ENG","UC Santa Cruz","Owns SDKs and developer experience, the one-line install that gets Revive into every agent stack.")}
  </div>
</div>""", 11))

# ---------- 12 VALIDATION ----------
S.append(frame(f"""
<div style="position:absolute;left:0;top:0;bottom:0;width:760px;background:linear-gradient(170deg,#1a1f2b,#12151d)">
  <div style="padding:110px 90px">
    <div class="mono" style="font-size:15px;letter-spacing:.42em;color:var(--lilac)">ALREADY DONE</div>
    <div style="margin-top:70px;display:flex;flex-direction:column;gap:42px">
      <div style="font-size:29px;font-weight:700;color:#fff"><span style="color:var(--green-br)">✓</span>&nbsp;&nbsp;Live product</div>
      <div style="font-size:29px;font-weight:700;color:#fff"><span style="color:var(--green-br)">✓</span>&nbsp;&nbsp;Works in production</div>
      <div style="font-size:29px;font-weight:700;color:#fff"><span style="color:var(--green-br)">✓</span>&nbsp;&nbsp;Developer tools shipped</div>
    </div>
    <div style="position:absolute;bottom:110px;left:90px;right:90px;font-size:21px;line-height:1.6;color:var(--slate)">
      <b style="color:#fff">Honest: no customers yet.</b><br>This is the 60-day plan to change that.</div>
  </div>
</div>
<div style="position:absolute;left:760px;right:0;top:0;bottom:0;padding:110px 120px 110px 110px">
  <div class="eyebrow">Next 60 days</div>
  <h1 style="font-size:58px;margin-top:26px">It's built. Now we prove<br>people want it.</h1>
  <div style="margin-top:76px;display:flex;flex-direction:column;gap:46px">
    <div style="border-left:6px solid var(--cobalt);padding-left:32px">
      <span class="mono" style="font-size:23px;color:var(--cobalt-deep);font-weight:700">30+</span>
      <span style="font-size:27px;font-weight:700;margin-left:16px">interviews with agent teams</span>
      <div style="font-size:19px;color:var(--muted);margin-top:8px">get the pain in their words</div></div>
    <div style="border-left:6px solid var(--cobalt);padding-left:32px">
      <span class="mono" style="font-size:23px;color:var(--cobalt-deep);font-weight:700">5</span>
      <span style="font-size:27px;font-weight:700;margin-left:16px">design-partner pilots</span>
      <div style="font-size:19px;color:var(--muted);margin-top:8px">free plan for a real case study</div></div>
    <div style="border-left:6px solid var(--cobalt);padding-left:32px">
      <span style="font-size:27px;font-weight:700">measure demand</span>
      <div style="font-size:19px;color:var(--muted);margin-top:8px">installs · waitlist · community</div></div>
    <div style="border-left:6px solid var(--cobalt);padding-left:32px">
      <span class="mono" style="font-size:23px;color:var(--cobalt-deep);font-weight:700">#1</span>
      <span style="font-size:27px;font-weight:700;margin-left:16px">publish the first pilot</span>
      <div style="font-size:19px;color:var(--muted);margin-top:8px">the first real numbers</div></div>
  </div>
</div>""", 12, ghost=False, runner=False, ticks=False))

# ---------- 13 ASK ----------
S.append(frame(f"""
<div style="position:absolute;left:0;top:0;bottom:0;width:22px;background:var(--cobalt)"></div>
<div class="pad z">
  <div class="eyebrow">The ask</div>
  <div style="font-family:'Space Grotesk';font-weight:700;font-size:290px;letter-spacing:-.05em;line-height:1;margin-top:30px">$750k</div>
  <div style="display:flex;align-items:center;gap:34px;margin-top:34px">
    <span class="mono" style="font-size:17px;letter-spacing:.4em;color:var(--cobalt-deep);font-weight:700">18-MONTH RUNWAY</span>
    <span style="font-size:23px;color:var(--muted)">one risk left: distribution. Here is what the money turns it into.</span>
  </div>
  <div style="display:flex;gap:80px;margin-top:80px">
    <div style="flex:1;border-top:4px solid var(--ink);padding-top:28px">
      <div style="font-family:'Space Grotesk';font-weight:700;font-size:62px;color:var(--cobalt-deep)">5 → 10+</div>
      <div style="font-size:22px;font-weight:600;margin-top:12px">pilots turned into paying teams</div></div>
    <div style="flex:1;border-top:4px solid var(--ink);padding-top:28px">
      <div style="font-family:'Space Grotesk';font-weight:700;font-size:62px;color:var(--cobalt-deep)">3 + 1</div>
      <div style="font-size:22px;font-weight:600;margin-top:12px">integrations live, first enterprise</div></div>
    <div style="flex:1;border-top:4px solid var(--ink);padding-top:28px">
      <div style="font-family:'Space Grotesk';font-weight:700;font-size:62px;color:var(--cobalt-deep)">~$450k</div>
      <div style="font-size:22px;font-weight:600;margin-top:12px">ARR, ready to raise the A</div></div>
  </div>
</div>
<div class="kicker"><span class="k">The money buys</span>
  <span style="font-size:23px;color:#fff">runway for the team · first sales hire · two more integrations · security groundwork</span>
</div>""", 13))

# ---------- 14 CLOSE ----------
S.append(frame(f"""
<div style="position:absolute;right:-60px;bottom:-140px;font-family:'Space Grotesk';font-weight:700;font-size:640px;color:#181d29;line-height:1">R.</div>
<div class="pad z">
  <div class="eyebrow">Revive Labs</div>
  <div style="font-family:'Space Grotesk';font-weight:700;font-size:98px;line-height:1.1;letter-spacing:-.03em;margin-top:130px">
    <span style="color:#fff">Agents are getting hands.</span><br>
    <span style="color:var(--lilac)">Someone has to own what<br>happens when the hands slip.</span>
  </div>
  <div style="position:absolute;bottom:112px;left:120px;font-size:25px">
    <b style="color:#fff">founders@revivelabs.app</b>
    <span style="color:var(--slate)">&nbsp;&nbsp;·&nbsp;&nbsp;revivelabs.app&nbsp;&nbsp;·&nbsp;&nbsp;deck and data room on request</span>
  </div>
</div>""", 14, dark=True, ghost=False, runner=False))

# ---- write html, screenshot, assemble ----
CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"
for i, html in enumerate(S, 1):
    path = os.path.join(SLIDES_DIR, f"s{i:02d}.html")
    with open(path, "w") as f: f.write(html)

import time
for i in range(1, len(S)+1):
    src = os.path.join(SLIDES_DIR, f"s{i:02d}.html")
    out = os.path.join(PNG_DIR, f"s{i:02d}.png")
    if os.path.exists(out): os.remove(out)
    # Chrome 149 headless writes the screenshot but never exits: launch, poll for file, kill.
    proc = subprocess.Popen([CHROME, "--headless=new", "--disable-gpu", "--hide-scrollbars",
                    "--no-sandbox", "--disable-crashpad", "--use-mock-keychain", "--password-store=basic",
                    "--no-first-run", f"--user-data-dir=/tmp/chrome-deck-{i}",
                    "--window-size=1920,1080", "--force-device-scale-factor=2",
                    f"--screenshot={out}", f"file://{src}"],
                   stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    deadline = time.time() + 45
    size = -1
    while time.time() < deadline:
        if os.path.exists(out):
            sz = os.path.getsize(out)
            if sz > 0 and sz == size: break   # stable size two polls in a row
            size = sz
        time.sleep(0.6)
    proc.kill(); proc.wait()
    subprocess.run(["rm","-rf",f"/tmp/chrome-deck-{i}"], capture_output=True)
    if not os.path.exists(out): raise RuntimeError(f"render failed for slide {i}")
    print("rendered", out)

from pptx import Presentation
from pptx.util import Inches
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
for i in range(1, len(S)+1):
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(os.path.join(PNG_DIR, f"s{i:02d}.png"), 0, 0, prs.slide_width, prs.slide_height)
    if i == 11:
        s.notes_slide.notes_text_frame.text = ("TO FINISH: drop headshots into decks/photos/ss.png, rg.png, ap.png and rerun "
            "python3 html_deck.py. Add one real credential per founder line. Titles: Srikanth CEO, Revanth CTO, Aarush Eng.")
out = os.path.join(HERE, "revive-seed-deck.pptx")
prs.save(out)
print("saved", out, os.path.getsize(out), "bytes,", len(S), "slides")
